import pandas as pd
from geopy.distance import geodesic
import streamlit as st
import folium
from streamlit_folium import st_folium
from pptx import Presentation
from pptx.util import Inches, Pt
import requests
import urllib.parse
from branca.element import Template, MacroElement
import os
import tempfile
import streamlit.components.v1 as components

# --- Streamlit page config ---
st.set_page_config(page_title="Closest Centres Map", layout="wide")

# --- Hide Streamlit UI elements ---
st.markdown("""
    <style>
    #MainMenu, footer, header, a[href*="github.com"], .viewerBadge_container__1QSob, .stAppViewerBadge {
        display: none !important;
    }
    </style>
""", unsafe_allow_html=True)

# --- JS to hide dynamically injected badges ---
components.html("""
<script>
setInterval(() => {
  const badges = document.querySelectorAll(
    '.viewerBadge_container__1QSob, .stAppViewerBadge, a[href*="github.com"]');
  badges.forEach(badge => badge.style.display = 'none');
}, 500);
</script>
""", height=0)

# --- Login system ---
def login():
    st.image("IWG Logo.jpg", width=150)
    st.title("Internal Map Login")
    email = st.text_input("Email")
    password = st.text_input("Password", type="password")
    if st.button("Login"):
        if password == "IWG123" and email.endswith("@iwgplc.com"):
            st.session_state["authenticated"] = True
            st.session_state["user_email"] = email
            st.success("Login successful!")
            st.experimental_rerun()
        else:
            st.error("Invalid email or password.")

if "authenticated" not in st.session_state:
    st.session_state["authenticated"] = False
if not st.session_state["authenticated"]:
    login()
    st.stop()

# --- Helper functions ---
def infer_area_type(location):
    formatted_str = location.get("formatted", "").lower()
    cbd_keywords = ["new york","manhattan","brooklyn","queens","bronx","staten island",
                    "chicago","los angeles","san francisco","boston","washington","philadelphia",
                    "houston","seattle","miami","atlanta","dallas","phoenix","detroit",
                    "san diego","minneapolis","denver","austin","portland","nashville",
                    "new orleans","las vegas","toronto","montreal","vancouver","ottawa",
                    "calgary","edmonton","winnipeg","halifax","victoria","quebec city",
                    "mexico city","guadalajara","monterrey","tijuana"]
    suburb_keywords = ["westmount","laval","longueuil","brossard","côte-saint-luc","ndg",
                      "saint-laurent","west island","mississauga","brampton","markham",
                      "vaughan","richmond hill","pickering","ajax","oshawa","milton",
                      "oakville","burlington","burnaby","surrey","richmond bc","coquitlam",
                      "delta","langley","maple ridge","north vancouver","west vancouver",
                      "okotoks","airdrie","sherwood park","st. albert","gatineau","kanata",
                      "orleans","cambridge","brookline","somerville","newton","quincy",
                      "jersey city","hoboken","newark","yonkers","staten island","flushing",
                      "long island city","bronxville","white plains","oakland","berkeley",
                      "san mateo","redwood city","palo alto","pasadena","burbank",
                      "santa monica","long beach","anaheim","evanston","oak park",
                      "naperville","schaumburg","coral gables","hialeah","kendall","aventura",
                      "zapopan","tlajomulco","santa catarina","san nicolas de los garza"]
    if any(city in formatted_str for city in cbd_keywords):
        return "CBD"
    elif any(nhood in formatted_str for nhood in suburb_keywords):
        return "Suburb"
    else:
        return "Suburb"

def normalize_centre_number(val):
    if pd.isna(val): return ""
    val_str = str(val).strip()
    return val_str.lstrip("0") or "0"

def normalize_address(addr):
    if pd.isna(addr): return ""
    return addr.strip().lower()

def filter_duplicates(df):
    preferred_statuses = {
        "Under Construction","Contract Signed","IC Approved",
        "Not Paid But Contract Signed","Centre Open"
    }
    df["Normalized Address"] = df["Addresses"].apply(normalize_address)
    grouped = df.groupby(["Centre Number","Normalized Address"])
    def select_preferred(group):
        preferred = group[group["Transaction Milestone Status"].isin(preferred_statuses)]
        return preferred if not preferred.empty else group
    filtered_df = grouped.apply(select_preferred).reset_index(drop=True)
    return filtered_df.drop(columns=["Normalized Address"])

# --- Main UI ---
st.title("\U0001F4CD Find 5 Closest Centres")
api_key = "edd4cb8a639240daa178b4c6321a60e6"
input_address = st.text_input("Enter an address:")

if input_address:
    try:
        with st.spinner("Loading, please wait..."):
            encoded_address = urllib.parse.quote(input_address)
            url = f"https://api.opencagedata.com/geocode/v1/json?q={encoded_address}&key={api_key}"
            response = requests.get(url)
            data = response.json()

            if response.status_code != 200:
                st.error(f"\u274C API Error: {response.status_code}. Try again.")
            elif not data.get('results'):
                st.error("\u274C Address not found. Try again.")
            else:
                location = data['results'][0]
                input_coords = (location['geometry']['lat'], location['geometry']['lng'])
                area_type = infer_area_type(location)
                st.write(f"Area type detected: **{area_type}**")

                file_path = "Database IC.xlsx"
                sheets = ["Comps", "Active Centre", "Centre Opened"]
                all_data = []
                for sheet in sheets:
                    df = pd.read_excel(file_path, sheet_name=sheet, engine="openpyxl")
                    df["Centre Number"] = df["Centre Number"].apply(normalize_centre_number)
                    if sheet in ["Active Centre", "Centre Opened"]:
                        df["Addresses"] = df["Address Line 1"]
                    else:
                        if "Addresses" not in df.columns and "Address Line 1" in df.columns:
                            df["Addresses"] = df["Address Line 1"]
                    df["Source Sheet"] = sheet
                    all_data.append(df)

                combined_data = pd.concat(all_data, ignore_index=True)
                combined_data = combined_data.dropna(subset=["Latitude","Longitude","Centre Number"])

                def has_valid_address(val):
                    return False if pd.isna(val) or (isinstance(val, str) and val.strip() == "") else True
                dupe_centre_nums = combined_data["Centre Number"][combined_data["Centre Number"].duplicated(keep=False)].unique()
                condition = combined_data["Centre Number"].isin(dupe_centre_nums) & (~combined_data["Addresses"].apply(has_valid_address))
                combined_data = combined_data[~condition]

                priority_order = {"Comps":0,"Active Centre":1,"Centre Opened":2}
                combined_data["Sheet Priority"] = combined_data["Source Sheet"].map(priority_order)
                data = combined_data.sort_values(by="Sheet Priority").drop_duplicates(subset=["Centre Number"],keep="first").drop(columns=["Sheet Priority"])

                active_centre_df = pd.read_excel(file_path, sheet_name="Active Centre", engine="openpyxl")
                active_centre_df["Centre Number"] = active_centre_df["Centre Number"].apply(normalize_centre_number)
                active_status_map = active_centre_df.dropna(subset=["Centre Number","Transaction Milestone Status"]).set_index("Centre Number")["Transaction Milestone Status"].to_dict()
                def replace_transaction_status(row):
                    return active_status_map[row["Centre Number"]] if row["Centre Number"] in active_status_map else row["Transaction Milestone Status"]
                data["Transaction Milestone Status"] = data.apply(replace_transaction_status, axis=1)
                data = filter_duplicates(data)
                for col in ["City","State","Zipcode"]:
                    if col not in data.columns: data[col] = ""
                data["Distance (miles)"] = data.apply(lambda row: geodesic(input_coords, (row["Latitude"], row["Longitude"])).miles, axis=1)
                data_sorted = data.sort_values("Distance (miles)").reset_index(drop=True)

                selected_centres, seen_distances, seen_centre_numbers = [], [], set()
                for _, row in data_sorted.iterrows():
                    d = row["Distance (miles)"]
                    centre_num = row["Centre Number"]
                    if centre_num in seen_centre_numbers:
                        continue
                    if all(abs(d - x) >= 0.005 for x in seen_distances):
                        selected_centres.append(row)
                        seen_centre_numbers.add(centre_num)
                        seen_distances.append(d)
                    if len(selected_centres) == 5:
                        break
                closest = pd.DataFrame(selected_centres)

                m = folium.Map(location=input_coords, zoom_start=14, zoom_control=True, control_scale=True)
                folium.Marker(location=input_coords, popup=f"Your Address: {input_address}", icon=folium.Icon(color="green")).add_to(m)

                def get_marker_color(ftype):
                    return {"Regus":"blue","HQ":"darkblue","Signature":"purple","Spaces":"black","Non-Standard Brand":"gold"}.get(ftype,"red")

                distance_text = ""
                for _, row in closest.iterrows():
                    dest_coords = (row["Latitude"], row["Longitude"])
                    folium.PolyLine([input_coords, dest_coords], color="blue", weight=2.5).add_to(m)
                    color = get_marker_color(row["Format - Type of Centre"])
                    label = f"#{int(row['Centre Number'])} - ({row['Distance (miles)']:.2f} mi)"
                    folium.Marker(location=dest_coords,
                                  popup=(f"#{int(row['Centre Number'])} - {row['Addresses']} | {row.get('City','')}, {row.get('State','')} {row.get('Zipcode','')} | "
                                         f"{row['Format - Type of Centre']} | {row['Transaction Milestone Status']} | {row['Distance (miles)']:.2f} mi"),
                                  tooltip=folium.Tooltip(f"<div style='font-size:16px;font-weight:bold'>{label}</div>", permanent=True, direction='right'),
                                  icon=folium.Icon(color=color)).add_to(m)
                    distance_text += f"Centre #{int(row['Centre Number'])} - {row['Addresses']}, {row.get('City','')}, {row.get('State','')} {row.get('Zipcode','')} - Format: {row['Format - Type of Centre']} - Milestone: {row['Transaction Milestone Status']} - {row['Distance (miles)']:.2f} miles\n"

                radius_miles = {"CBD":1,"Suburb":5,"Rural":10}
                radius_m = radius_miles.get(area_type,5) * 1609.34
                folium.Circle(location=input_coords, radius=radius_m, color="green", fill=True, fill_opacity=0.2).add_to(m)

                # Patched legend for radius (text visible)
                legend_html = f"""
                    {{% macro html(this, kwargs) %}}
                    <div style='position: absolute; top: 70px; left: 10px; width: 180px; z-index: 9999;
                                background-color: white; padding: 10px; border: 2px solid gray;
                                border-radius: 5px; font-size: 14px; color: black; text-shadow: 1px 1px 2px white;'>
                        <b>Radius</b><br>
                        <span style='color:green;'>&#x25CF;</span> {radius_miles.get(area_type,5)}-mile Zone
                    </div>
                    {{% endmacro %}}
                """
                legend = MacroElement()
                legend._template = Template(legend_html)
                m.get_root().add_child(legend)

                col1, col2 = st.columns([5, 2])
                with col1:
                    st_folium(m, width=950, height=650)
                    st.markdown(f"""
                        <div style="font-size:18px; line-height:1.5; font-weight:bold; padding-top: 8px;">
                        {distance_text.replace(chr(10), "<br>")}
                        </div>
                    """, unsafe_allow_html=True)
                with col2:
                    # Patched legend for centre types with visible black text and white text-shadow
                    st.markdown("""
                        <div style="
                            background-color: white;
                            padding: 10px;
                            border: 2px solid grey;
                            border-radius: 10px;
                            width: 100%;
                            margin-top: 20px;
                            color: black;
                            text-shadow: 1px 1px 2px white;
                            font-weight: bold;
                            font-size: 14px;
                        ">
                            Centre Type Legend<br>
                            <i style="background-color: lightgreen; padding: 5px;">&#9724;</i> Proposed Address<br>
                            <i style="background-color: lightblue; padding: 5px;">&#9724;</i> Regus<br>
                            <i style="background-color: darkblue; padding: 5px;">&#9724;</i> HQ<br>
                            <i style="background-color: purple; padding: 5px;">&#9724;</i> Signature<br>
                            <i style="background-color: black; padding: 5px;">&#9724;</i> Spaces<br>
                            <i style="background-color: gold; padding: 5px;">&#9724;</i> Non-Standard Brand
                        </div>
                    """, unsafe_allow_html=True)

                uploaded_image = st.file_uploader("📼 Optional: Upload Map Screenshot for PowerPoint", type=["png","jpg","jpeg"])

                if st.button("📤 Export to PowerPoint"):
                    try:
                        prs = Presentation()
                        slide_layout = prs.slide_layouts[5]
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = f"5 Closest Centres to:\n{input_address}"

                        if uploaded_image:
                            image_path = os.path.join(tempfile.gettempdir(), uploaded_image.name)
                            with open(image_path, "wb") as img_file:
                                img_file.write(uploaded_image.read())
                            slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), height=Inches(3.5))

                        def add_centres_to_slide_table(centres_subset, title=None):
                            slide = prs.slides.add_slide(slide_layout)
                            if title:
                                slide.shapes.title.text = title

                            rows = len(centres_subset) + 1
                            cols = 6
                            table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1), Inches(9), Inches(0.8 + 0.4 * rows)).table

                            headers = ["Centre #", "Address", "City, State, Zip", "Format", "Milestone", "Distance (miles)"]
                            for col_idx, header_text in enumerate(headers):
                                cell = table.cell(0, col_idx)
                                cell.text = header_text
                                for p in cell.text_frame.paragraphs:
                                    p.font.bold = True
                                    p.font.size = Pt(14)

                            for i, row in enumerate(centres_subset.itertuples(), start=1):
                                table.cell(i, 0).text = str(int(row._asdict()["Centre Number"]))
                                table.cell(i, 1).text = row._asdict()["Addresses"]
                                city = row._asdict().get("City", "")
                                state = row._asdict().get("State", "")
                                zipc = row._asdict().get("Zipcode", "")
                                table.cell(i, 2).text = f"{city}, {state} {zipc}".strip(", ")
                                table.cell(i, 3).text = row._asdict()["Format - Type of Centre"]
                                table.cell(i, 4).text = row._asdict()["Transaction Milestone Status"]
                                table.cell(i, 5).text = f"{row._asdict()['Distance (miles)']:.2f}"

                        add_centres_to_slide_table(closest, title="Closest Centres")

                        pptx_path = os.path.join(tempfile.gettempdir(), "ClosestCentres.pptx")
                        prs.save(pptx_path)
                        with open(pptx_path, "rb") as f:
                            pptx_bytes = f.read()
                        st.success("PowerPoint file created!")
                        st.download_button("⬇️ Download PowerPoint", pptx_bytes, file_name="ClosestCentres.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

                    except Exception as e:
                        st.error(f"Error generating PowerPoint: {e}")

    except Exception as ex:
        st.error(f"Unexpected error: {ex}")

